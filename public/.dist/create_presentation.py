from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = "Smart_Hostel_Network_Manager_Presentation.pptx"
NOTES_OUT = "Smart_Hostel_Network_Manager_Speaker_Notes.md"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(8, 13, 26)
PANEL = RGBColor(18, 31, 49)
PANEL_2 = RGBColor(22, 41, 64)
CYAN = RGBColor(18, 214, 231)
BLUE = RGBColor(86, 130, 205)
GREEN = RGBColor(0, 214, 118)
RED = RGBColor(255, 55, 91)
YELLOW = RGBColor(255, 205, 77)
TEXT = RGBColor(238, 244, 255)
MUTED = RGBColor(142, 165, 205)
WHITE = RGBColor(255, 255, 255)


def blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def add_text(slide, text, x, y, w, h, size=24, color=TEXT, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.65, 0.45, 11.8, 0.55, 30, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.18), Inches(1.4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.32, 11.5, 0.45, 15, MUTED)


def add_card(slide, x, y, w, h, title, body=None, accent=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL_2
    shape.line.width = Pt(1.2)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + 0.22, y + 0.16, w - 0.35, 0.35, 15, WHITE, True)
    if body:
        add_text(slide, body, x + 0.22, y + 0.64, w - 0.35, h - 0.8, 12.5, MUTED)
    return shape


def add_bullets(slide, items, x, y, w, h, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT if idx == 0 else MUTED
        p.space_after = Pt(8)
    return box


def stat(slide, x, y, label, value, color):
    add_card(slide, x, y, 2.25, 1.2, label, None, color)
    add_text(slide, value, x + 0.25, y + 0.5, 1.7, 0.45, 27, color, True)


def chip(slide, x, y, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_2
    shape.line.color.rgb = color
    add_text(slide, text, x + 0.12, y + 0.08, 1.55, 0.18, 9.5, color, True, PP_ALIGN.CENTER)


slides_notes = []

# Slide 1
s = blank_slide()
add_text(s, "Smart Hostel", 0.75, 1.0, 6.0, 0.5, 28, CYAN, True)
add_text(s, "Network Manager", 0.75, 1.52, 8.4, 0.75, 42, WHITE, True)
add_text(s, "A web-based system to register hostel devices, manage blocked devices, and monitor network usage.", 0.78, 2.45, 7.2, 0.8, 18, MUTED)
stat(s, 0.78, 4.1, "Frontend", "HTML/CSS/JS", CYAN)
stat(s, 3.25, 4.1, "Backend", "Node + Express", GREEN)
stat(s, 5.72, 4.1, "Database", "SQLite", YELLOW)
add_card(s, 9.0, 0.95, 3.35, 4.85, "Project Snapshot", "Dashboard, device registration, block/unblock controls, usage logs, and reports are integrated into a single local web application.", BLUE)
slides_notes.append(("Title", "Introduce the project as a hostel network administration tool for managing student devices and network usage."))

# Slide 2
s = blank_slide()
add_title(s, "Problem Statement", "Hostel networks need simple device control and visibility.")
add_bullets(s, [
    "Manual device tracking becomes difficult as student device count grows.",
    "Unauthorized or problematic devices need quick blocking and unblocking.",
    "Administrators need basic visibility into usage and active devices.",
    "A lightweight local system is easier to deploy for small hostels or labs."
], 0.9, 2.0, 5.8, 3.8)
add_card(s, 7.3, 2.0, 4.7, 2.8, "Goal", "Build a browser-based network manager that stores device records, controls blocked devices, and summarizes usage with a simple dashboard.", CYAN)
slides_notes.append(("Problem Statement", "Explain why the system is useful: it replaces scattered records with a central network management interface."))

# Slide 3
s = blank_slide()
add_title(s, "Objectives")
cards = [
    ("Register Devices", "Store MAC address, owner, room number, and device name.", CYAN),
    ("Control Access", "Block, unblock, and delete devices from the admin interface.", RED),
    ("Track Usage", "Log sent and received bytes for each registered device.", GREEN),
    ("Generate Reports", "Display total data usage, sessions, and last active time.", YELLOW),
]
for i, (t, b, c) in enumerate(cards):
    add_card(s, 0.85 + (i % 2) * 6.1, 1.75 + (i // 2) * 2.1, 5.35, 1.55, t, b, c)
slides_notes.append(("Objectives", "Walk through the four major objectives: registration, access control, usage tracking, and reports."))

# Slide 4
s = blank_slide()
add_title(s, "System Architecture", "Three simple layers keep the project easy to run and maintain.")
add_card(s, 0.75, 2.2, 3.1, 1.25, "Frontend", "index.html\nstyle.css\napp.js", CYAN)
add_card(s, 5.05, 2.2, 3.1, 1.25, "Backend API", "Express server\nREST endpoints\nValidation", GREEN)
add_card(s, 9.35, 2.2, 3.1, 1.25, "Database", "SQLite file\nbetter-sqlite3\nPersistent data", YELLOW)
for x in [4.05, 8.35]:
    line = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.58), Inches(0.7), Inches(0.35))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
add_text(s, "Browser UI sends fetch requests to /api routes. The server validates input and updates SQLite tables.", 1.1, 4.45, 11.2, 0.7, 18, MUTED, False, PP_ALIGN.CENTER)
slides_notes.append(("System Architecture", "Describe how the frontend, backend, and database interact through API requests."))

# Slide 5
s = blank_slide()
add_title(s, "Core Features")
add_bullets(s, [
    "Dashboard with total devices, active devices, blocked devices, pending devices, and total data.",
    "Device registration using MAC address, owner name, device name, and room number.",
    "All Devices table with add usage, block, unblock, and delete actions.",
    "Blocked Devices view for reviewing and restoring blocked MAC addresses.",
    "Usage Reports table showing sent data, received data, sessions, and last active time."
], 0.9, 1.75, 7.0, 4.6, 16)
chip(s, 8.6, 1.95, "Dashboard", CYAN)
chip(s, 8.6, 2.55, "Register", GREEN)
chip(s, 8.6, 3.15, "Block/Unblock", RED)
chip(s, 8.6, 3.75, "Reports", YELLOW)
slides_notes.append(("Core Features", "Mention the user-facing screens and the actions available to the admin."))

# Slide 6
s = blank_slide()
add_title(s, "Database Design")
add_card(s, 0.75, 1.65, 3.75, 3.9, "devices", "id\nmac_address\ndevice_name\nowner_name\nroom_number\nstatus\nregistered_at\nlast_seen", CYAN)
add_card(s, 4.8, 1.65, 3.75, 3.9, "usage_logs", "id\ndevice_id\nmac_address\nbytes_sent\nbytes_received\nlogged_at", GREEN)
add_card(s, 8.85, 1.65, 3.75, 3.9, "blocked_devices", "id\nmac_address\nreason\nblocked_at", RED)
add_text(s, "SQLite keeps the project lightweight while still supporting structured queries and persistent storage.", 1.15, 6.1, 11.0, 0.45, 16, MUTED, False, PP_ALIGN.CENTER)
slides_notes.append(("Database Design", "Explain the three tables and how MAC address connects device records with usage and blocking."))

# Slide 7
s = blank_slide()
add_title(s, "API Endpoints")
rows = [
    ("GET", "/api/devices", "List registered devices"),
    ("POST", "/api/devices/register", "Register a new device"),
    ("POST", "/api/block", "Block a MAC address"),
    ("DELETE", "/api/block/:mac", "Unblock a device"),
    ("POST", "/api/delete-device", "Delete device and related records"),
    ("POST", "/api/usage/log", "Add usage log"),
    ("GET", "/api/usage/report", "View usage report"),
    ("GET", "/api/stats", "Dashboard statistics"),
]
y = 1.55
for method, path, purpose in rows:
    add_text(s, method, 0.85, y, 1.1, 0.25, 12, CYAN, True)
    add_text(s, path, 2.0, y, 3.5, 0.25, 12, WHITE, True)
    add_text(s, purpose, 5.7, y, 6.2, 0.25, 12, MUTED)
    y += 0.55
slides_notes.append(("API Endpoints", "Show that the backend exposes clear REST-style routes for frontend operations."))

# Slide 8
s = blank_slide()
add_title(s, "Workflow")
steps = [
    ("1", "Admin opens dashboard"),
    ("2", "Registers student device"),
    ("3", "Device appears in All Devices"),
    ("4", "Admin can add usage, block, unblock, or delete"),
    ("5", "Reports summarize network activity"),
]
x_positions = [0.7, 3.2, 5.7, 8.2, 10.7]
for (num, label), x in zip(steps, x_positions):
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.0), Inches(0.75), Inches(0.75))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CYAN
    circle.line.fill.background()
    add_text(s, num, x + 0.28, 2.17, 0.2, 0.2, 17, BG, True)
    add_text(s, label, x - 0.35, 3.0, 1.55, 0.9, 12, MUTED, False, PP_ALIGN.CENTER)
for x in [1.65, 4.15, 6.65, 9.15]:
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.18), Inches(1.1), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.fill.background()
slides_notes.append(("Workflow", "Use this slide to explain a typical admin flow from dashboard to reporting."))

# Slide 9
s = blank_slide()
add_title(s, "Implementation Highlights")
add_bullets(s, [
    "Express serves both the frontend files and JSON API routes on port 3001.",
    "better-sqlite3 provides simple synchronous database access for local deployment.",
    "Device deletion removes related usage and blocked-device records to keep data consistent.",
    "Frontend uses fetch requests and refreshes the active views after each action.",
    "The application can be run with npm start from the backend folder."
], 0.9, 1.75, 6.8, 4.6, 16)
add_card(s, 8.2, 2.05, 3.9, 2.2, "Run Command", "cd backend\nnpm start\n\nOpen:\nhttp://localhost:3001", GREEN)
slides_notes.append(("Implementation Highlights", "Summarize the technical choices and mention how to run the project."))

# Slide 10
s = blank_slide()
add_title(s, "Advantages And Future Scope")
add_card(s, 0.85, 1.8, 5.4, 3.5, "Advantages", "Simple interface\nLocal database\nQuick device control\nUsage visibility\nEasy to demo and extend", CYAN)
add_card(s, 7.0, 1.8, 5.4, 3.5, "Future Scope", "Admin login\nReal router integration\nCharts and filters\nExport reports\nRole-based permissions", YELLOW)
slides_notes.append(("Advantages And Future Scope", "Close by explaining the current value and how the project can be improved later."))

# Slide 11
s = blank_slide()
add_text(s, "Thank You", 0.9, 2.0, 11.5, 0.7, 48, WHITE, True, PP_ALIGN.CENTER)
add_text(s, "Smart Hostel Network Manager", 0.9, 3.05, 11.5, 0.4, 22, CYAN, True, PP_ALIGN.CENTER)
add_text(s, "Questions?", 0.9, 3.75, 11.5, 0.35, 20, MUTED, False, PP_ALIGN.CENTER)
slides_notes.append(("Thank You", "End the presentation and invite questions."))

prs.save(OUT)

with open(NOTES_OUT, "w", encoding="utf-8") as f:
    f.write("# Smart Hostel Network Manager - Speaker Notes\n\n")
    for i, (title, note) in enumerate(slides_notes, start=1):
        f.write(f"## Slide {i}: {title}\n{note}\n\n")

print(OUT)
print(NOTES_OUT)
